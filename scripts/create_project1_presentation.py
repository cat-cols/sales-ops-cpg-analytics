from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project 1: Ops Command Center"
    subtitle.text = "End-to-End Data Flow Architecture"
    
    # Slide 2: Executive Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Executive Overview"
    content.text = """Project 1 demonstrates end-to-end analytics engineering:

• Ingests messy cross-functional source data
• Standardizes and models data into clean, reliable tables
• Reconciles outputs back to finance-facing totals
• Builds QA controls with hard-fail mechanisms
• Prepares data foundation for BI dashboards
• Translates data into actionable business insights

This is the flagship project showcasing Business Analyst / Analytics Engineering workflow."""
    
    # Slide 3: Architecture Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Architecture Overview"
    content.text = """Layered Data Architecture:

Source Systems → Raw → Staging → Integration → Mart → QA → BI

• Source: Sales, Operations, Labor, Finance, Reference data
• Raw: PostgreSQL raw tables (as-loaded from sources)
• Staging: Cleaned and typed views (STG layer)
• Integration: Conformed views (INT layer)
• Mart: Business-ready facts, dimensions, KPIs (MART layer)
• QA: Quality assurance controls and reconciliation
• BI: Tableau dashboards for executive reporting"""
    
    # Slide 4: Source Systems
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Source Systems"
    content.text = """Cross-Functional Data Sources:

• Sales Data: POS transactions, distributor shipments
• Operations/Inventory: Stock levels, shipments, receipts
• Labor: Hours, headcount, costs, productivity metrics
• Finance: GL actuals, budget targets
• Reference: Product catalog, store locations, employee data

Intentional Challenges:
• Duplicate rows and missing keys
• Inconsistent formats and data types
• Cross-system reconciliation requirements
• Realistic data quality issues"""
    
    # Slide 5: Data Ingestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Data Ingestion"
    content.text = """Source to Raw Layer:

Process:
1. Generate synthetic source data (Python scripts)
2. Load into PostgreSQL raw tables
3. Preserve source system structure
4. Enable data lineage tracking

Key Features:
• Truncate-then-append for reproducibility
• Source register for tracking
• Data freshness monitoring
• Row count validation"""
    
    # Slide 6: Staging Layer (STG)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Staging Layer (STG)"
    content.text = """Purpose: Clean and standardize source data

Key Transformations:
• Data type standardization (strings to dates, numbers)
• Null handling and default values
• Column naming conventions
• Duplicate removal
• Business rule validation

Output:
• Clean, typed views ready for integration
• Data quality flags
• Source system lineage preserved
• 20+ staging views across all domains"""
    
    # Slide 7: Integration Layer (INT)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Integration Layer (INT)"
    content.text = """Purpose: Conform data across sources

Key Functions:
• Truth selection (POS vs. distributor for sales)
• Cross-system conformance
• Canonical dimensions (store, SKU, date)
• Grain alignment
• Business logic application

Output:
• Conformed facts (sales, inventory, labor)
• Shared dimensions (store, SKU, date, employee)
• Single source of truth for each entity
• 15+ integration views"""
    
    # Slide 8: Mart Layer (MART)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Mart Layer (MART)"
    content.text = """Purpose: Business-ready data for reporting

Components:
• Facts: Sales, inventory, labor, finance actuals
• Dimensions: Store, SKU, date, employee
• KPIs: Gross margin, labor productivity, in-stock rate
• Controls: Reconciliation, data quality
• Recon views: Sales to GL, source-to-source

Key Features:
• Star schema design
• Pre-calculated KPIs
• Business-friendly naming
• 25+ mart views"""
    
    # Slide 9: KPI Layer
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "KPI Layer"
    content.text = """Pre-Calculated Business Metrics:

Sales KPIs:
• Net sales, gross sales, discount rate
• Average order value, units per order
• Customer count, order count

Profitability KPIs:
• Gross margin dollars and percentage
• COGS, discount impact

Labor KPIs:
• Sales per labor hour
• Labor cost percentage
• Overtime rate

Inventory KPIs:
• In-stock rate, fill rate
• Days of supply, backorder rate

All KPIs documented in metric_dictionary.md"""
    
    # Slide 10: QA & Reconciliation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "QA & Reconciliation"
    content.text = """Quality Assurance:

Hard-Fail Controls:
• Row count monitoring
• Missing dimension join checks
• Grain/uniqueness validation
• Freshness monitoring
• Mart dimension contract checks

Reconciliation:
• Sales to GL variance (1% tolerance)
• Distributor vs. POS comparison
• Finance tie-out logic
• Explicit tolerance rules

Output:
• Pass/fail status for all checks
• Variance explanations
• Trust indicators for BI layer"""
    
    # Slide 11: BI Layer (Tableau)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "BI Layer (Tableau)"
    content.text = """Dashboard Implementation:

Data Source:
• PostgreSQL connection to mart views
• Relationship model (fact + dimensions)
• Extracts for performance

Planned Dashboards:
• Executive Overview (KPI cards, trends)
• Sales & Margin Analysis
• Inventory & Operations
• Labor & Productivity
• Reconciliation & Data Health

Key Features:
• Cross-platform (works on Mac)
• LOD expressions for advanced analytics
• Parameter-driven analysis
• Published to Tableau Public"""
    
    # Slide 12: Business Value
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Business Value Demonstration"
    content.text = """Beyond Technical Implementation:

Store Performance Analysis:
• Identified 3 underperforming stores
• Root cause analysis with specific findings
• Remediation plans with financial projections
• $44.5K monthly recovery projected

Variance Analysis:
• Sales variance decomposition (+$127K)
• Margin variance drivers identified
• Beyond simple flagging to root cause

Strategic Recommendations:
• Channel strategy optimization (+$45K)
• Inventory rebalancing (+$20K)
• Labor productivity improvements (+$38K)

Total Projected Impact: $245K monthly improvement

See: docs/business_analysis_examples.md"""
    
    # Slide 13: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Technology Stack"
    content.text = """Data Warehouse:
• PostgreSQL (wyld_chyld database)
• 79 views across all layers
• SchemaSpy documentation

Data Processing:
• Python (pandas, numpy, pyarrow)
• SQL for data modeling
• psycopg for PostgreSQL connection

Data Quality:
• Great Expectations
• Pandera validation

BI & Visualization:
• Tableau (primary implementation)
• Power BI (semantic model planning)

Documentation:
• Markdown documentation
• Architecture diagrams
• Metric dictionaries"""
    
    # Slide 14: Key Files & Documentation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Key Files & Documentation"
    content.text = """Implementation:
• 01_ops_command_center/sql/ - SQL views by layer
• 01_ops_command_center/docs/ - Project documentation
• 01_ops_command_center/tableau/ - BI implementation

Documentation:
• executive_walkthrough.md - High-level overview
• executive_narrative.md - Detailed business context
• metric_dictionary.md - Metric definitions and formulas
• business_analysis_examples.md - Business value demonstration

BI Documentation:
• tableau/IMPLEMENTATION_GUIDE.md - Dashboard planning
• tableau/DATA_SOURCE_GUIDE.md - Connection setup
• tableau/CALCULATED_FIELDS.md - Calculation reference
• tableau/TABLEAU_VS_POWERBI.md - Platform comparison"""
    
    # Slide 15: How to Run
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "How to Run"
    content.text = """SQL Pipeline:
```bash
python3 scripts/generate_project1_data.py --pg-dsn "$P1_PG_OPS" --pg-load-mode truncate_then_append && \\
psql "$P1_PG_OPS" -v ON_ERROR_STOP=1 -f 01_ops_command_center/sql/stg/_build_stg.sql && \\
psql "$P1_PG_OPS" -v ON_ERROR_STOP=1 -f 01_ops_command_center/sql/int/_build_int.sql && \\
psql "$P1_PG_OPS" -v ON_ERROR_STOP=1 -f 01_ops_command_center/sql/mart/_build_mart.sql && \\
psql "$P1_PG_OPS" -v ON_ERROR_STOP=1 -f 01_ops_command_center/sql/_qa/_run_qa.sql
```

Tableau Dashboard:
1. Install Tableau Desktop (trial available)
2. Follow tableau/DATA_SOURCE_GUIDE.md
3. Build dashboards per tableau/IMPLEMENTATION_GUIDE.md
4. Publish to Tableau Public"""
    
    # Slide 16: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Summary"
    content.text = """Project 1 Demonstrates:

Technical Excellence:
• End-to-end data pipeline (source to BI)
• Layered SQL architecture (79 views)
• QA controls with hard-fail mechanisms
• Reconciliation to finance totals

Business Value:
• Store performance analysis with recommendations
• Variance decomposition beyond flagging
• $245K monthly improvement projected
• Actionable insights from data

Portfolio Value:
• Cross-platform BI implementation (Tableau)
• Platform versatility (Tableau + Power BI planning)
• Strong documentation and communication
• Realistic business scenario simulation

This is the flagship project showcasing Business Analyst / Analytics Engineering workflow."""
    
    # Save presentation
    prs.save('/Users/b/data/projects/althea-sales-ops-cpg/docs/project1_end_to_end.pptx')
    print("Created PowerPoint: docs/project1_end_to_end.pptx")

if __name__ == "__main__":
    create_presentation()
